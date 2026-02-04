# /backend/app/services/document_processor.py

from docx import Document
import openpyxl
from pptx import Presentation
from typing import Dict, List
import os

class DocumentProcessor:
    """Service for processing DOCX, XLSX, PPTX files"""
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += "\n" + cell.text
            
            return text.strip()
        except Exception as e:
            print(f"DOCX extraction error: {e}")
            return ""
    
    @staticmethod
    def extract_data_from_excel(file_path: str) -> Dict:
        """Extract data from Excel file"""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            sheets_data = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Get all rows
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out empty rows
                    if any(cell is not None for cell in row):
                        rows.append([str(cell) if cell is not None else "" for cell in row])
                
                if rows:
                    sheets_data[sheet_name] = {
                        "headers": rows[0] if len(rows) > 0 else [],
                        "rows": rows[1:] if len(rows) > 1 else [],
                        "row_count": len(rows) - 1,
                        "column_count": len(rows[0]) if rows else 0
                    }
            
            return sheets_data
        
        except Exception as e:
            print(f"Excel extraction error: {e}")
            return {}
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n=== Slide {slide_num} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"PPTX extraction error: {e}")
            return ""
    
    @staticmethod
    def extract_tables_from_docx(file_path: str) -> List[Dict]:
        """Extract tables from DOCX file"""
        try:
            doc = Document(file_path)
            tables_data = []
            
            for table_num, table in enumerate(doc.tables, 1):
                rows = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    rows.append(row_data)
                
                if rows:
                    table_dict = {
                        "table_number": table_num,
                        "headers": rows[0] if len(rows) > 0 else [],
                        "rows": rows[1:] if len(rows) > 1 else [],
                        "row_count": len(rows) - 1,
                        "column_count": len(rows[0]) if rows else 0
                    }
                    tables_data.append(table_dict)
            
            return tables_data
        
        except Exception as e:
            print(f"DOCX table extraction error: {e}")
            return []

document_processor = DocumentProcessor()